import json
import re
from pathlib import Path

from docx import Document
from google.genai import types
from pptx import Presentation

from config.gemini_client import client

FILE_TYPE_PATTERNS = {
    "pptx": re.compile(
        r"\b(pptx?|powerpoint|presentation|slide\s*deck|slides?\s*file)\b",
        re.IGNORECASE,
    ),
    "docx": re.compile(
        r"\b(docx?|word\s+doc(?:ument)?|microsoft\s+word|word\s+file)\b",
        re.IGNORECASE,
    ),
}

def detect_requested_file_type(text: str) -> str | None:
    for file_type, pattern in FILE_TYPE_PATTERNS.items():
        if pattern.search(text):
            return file_type
    return None


def _safe_filename(title: str, fallback: str) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*]', "_", title.strip())
    return cleaned or fallback


def build_pptx(slides_data: dict, output_path: Path) -> None:
    presentation = Presentation()

    title_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = slides_data.get("title", "Presentation")

    subtitle = slides_data.get("subtitle", "")
    if subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle

    content_layout = presentation.slide_layouts[1]
    for slide_data in slides_data.get("slides", []):
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get("title", "Slide")

        bullets = slide_data.get("bullets", [])
        if not bullets:
            continue

        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()

        for index, bullet in enumerate(bullets):
            if index == 0:
                text_frame.text = str(bullet)
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0

    presentation.save(output_path)


def generate_pptx(task_description: str, docs_dir: Path, prefix: str) -> tuple[str, str]:
    """Ask Gemini for slide structure, then build a .pptx file locally."""
    prompt = f"""
Create a PowerPoint presentation for this task:

{task_description}

Return ONLY valid JSON with this structure:
{{
  "title": "Presentation title",
  "subtitle": "Optional subtitle",
  "slides": [
    {{
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2", "Point 3"]
    }}
  ]
}}

Rules:
1. Include a title slide worth of content in "title" and "subtitle".
2. Provide at least 5 content slides unless the topic is very small.
3. Keep bullets concise and presentation-ready.
4. No markdown. No explanations outside JSON.
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            response_mime_type="application/json"
        ),
    )

    slides_data = json.loads(response.text)
    docs_dir.mkdir(exist_ok=True)

    safe_title = _safe_filename(slides_data.get("title", ""), "presentation")
    output_path = docs_dir / f"{prefix}_{safe_title}.pptx"
    build_pptx(slides_data, output_path)

    summary = (
        f"Generated PowerPoint: {slides_data.get('title', 'Presentation')} "
        f"({len(slides_data.get('slides', []))} slides)"
    )
    return summary, str(output_path.resolve())


def build_docx(document_data: dict, output_path: Path) -> None:
    document = Document()
    document.add_heading(document_data.get("title", "Document"), level=0)

    for section in document_data.get("sections", []):
        heading = section.get("heading")
        if heading:
            document.add_heading(heading, level=1)

        for paragraph in section.get("paragraphs", []):
            document.add_paragraph(str(paragraph))

        for bullet in section.get("bullets", []):
            document.add_paragraph(str(bullet), style="List Bullet")

    document.save(output_path)


def generate_docx(task_description: str, docs_dir: Path, prefix: str) -> tuple[str, str]:
    """Ask Gemini for document structure, then build a .docx file locally."""
    prompt = f"""
Create a Word document for this task:

{task_description}

Return ONLY valid JSON with this structure:
{{
  "title": "Document title",
  "sections": [
    {{
      "heading": "Section heading",
      "paragraphs": ["Paragraph text"],
      "bullets": ["Bullet point 1", "Bullet point 2"]
    }}
  ]
}}

Rules:
1. Provide a clear document title.
2. Include multiple well-structured sections.
3. Use paragraphs for explanatory text and bullets for lists.
4. No markdown. No explanations outside JSON.
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            response_mime_type="application/json"
        ),
    )

    document_data = json.loads(response.text)
    docs_dir.mkdir(exist_ok=True)

    safe_title = _safe_filename(document_data.get("title", ""), "document")
    output_path = docs_dir / f"{prefix}_{safe_title}.docx"

    build_docx(document_data, output_path)

    summary = (
        f"Generated Word document: {document_data.get('title', 'Document')} "
        f"({len(document_data.get('sections', []))} sections)"
    )
    return summary, str(output_path.resolve())


FILE_GENERATORS = {
    "pptx": generate_pptx,
    "docx": generate_docx,
}